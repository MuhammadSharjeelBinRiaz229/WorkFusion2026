#!/usr/bin/env python3
"""Update WorkFusion FYP presentation into a premium, professional slide deck with perfect grid alignment and custom vector diagrams."""

import copy
import os
import shutil
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(ROOT, "assets")
SRC = os.path.join(ROOT, "Presentation_backup.pptx")
OUT = os.path.join(ROOT, "Presentation.pptx")


def populate_slide_title(slide, title_text, font_size_pt=34):
    """Safely populate slide title placeholder while keeping its alignment."""
    title_ph = slide.shapes.title
    if title_ph:
        tf = title_ph.text_frame
        p = tf.paragraphs[0]
        alignment = p.alignment
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = alignment
        r = p.add_run()
        r.text = title_text
        r.font.bold = True
        r.font.name = "Segoe UI"
        r.font.size = Pt(font_size_pt)


def add_premium_bullets(slide, bullets, left, top, width, height, font_size_pt=14, is_sidebar_slide=False):
    """Create a pixel-perfect, custom text box with premium, bold-prefix bullet points."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_bottom = Inches(0)
    
    # In sidebar (Style B) slides, vertical space is larger, so we can afford more space between bullets
    space_after_pt = 12 if is_sidebar_slide else 8
    
    for idx, bullet_text in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(space_after_pt)
        p.line_spacing = 1.15
        
        # Bolding label prefix if there is a colon (bold-prefix layout)
        if ":" in bullet_text and not bullet_text.startswith("http"):
            parts = bullet_text.split(":", 1)
            label = parts[0].strip() + ": "
            value = parts[1].strip()
            
            # Bold Run for label
            r1 = p.add_run()
            r1.text = "• " + label
            r1.font.bold = True
            r1.font.size = Pt(font_size_pt)
            r1.font.name = "Segoe UI"
            
            # Normal Run for value
            r2 = p.add_run()
            r2.text = value
            r2.font.bold = False
            r2.font.size = Pt(font_size_pt)
            r2.font.name = "Segoe UI"
        else:
            r = p.add_run()
            r.text = "• " + bullet_text
            r.font.bold = False
            r.font.size = Pt(font_size_pt)
            r.font.name = "Segoe UI"
            
    return txbox


def add_screenshot_with_aspect_ratio(slide, screenshot_path, left, top, max_width, max_height):
    """Add a screenshot to the slide, maintaining its aspect ratio, fitting within max_width and max_height, 
    and centering it horizontally within the allocated column."""
    if not os.path.exists(screenshot_path):
        return None
        
    try:
        with Image.open(screenshot_path) as img:
            orig_w, orig_h = img.size
        aspect_ratio = orig_w / orig_h
        
        # Start by assuming we use the full allowed width
        width = max_width
        height = width / aspect_ratio
        
        # If height exceeds the allowed max height, shrink width & height to fit height bounds
        if height > max_height:
            height = max_height
            width = height * aspect_ratio
            
        # Center horizontally within the original max_width bounding box
        horizontal_offset = (max_width - width) / 2
        centered_left = left + horizontal_offset
        
        pic = slide.shapes.add_picture(screenshot_path, centered_left, top, width=width)
        return pic
    except Exception as e:
        print(f"Error adding screenshot {screenshot_path}: {e}")
        return None


def format_table(table, font_size_pt=11):
    """Format cells of a PowerPoint table cleanly."""
    for r_idx, row in enumerate(table.rows):
        for c_idx, cell in enumerate(row.cells):
            text = cell.text.strip()
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = text
            r.font.name = "Segoe UI"
            r.font.size = Pt(font_size_pt)
            if r_idx == 0:
                r.font.bold = True
                p.alignment = PP_ALIGN.CENTER


def clone_slide_layout_shapes(source_slide, dest_slide):
    """Clones non-placeholder background shapes, layout lines and borders, excluding text boxes and warning icons."""
    for shape in source_slide.shapes:
        if shape.is_placeholder:
            continue
        # Skip copying custom content text boxes or shapes that contain actual text content
        if shape.has_text_frame and shape.text_frame.text.strip():
            continue
        # Also skip if it's a text box shape type (TEXT_BOX is 17)
        if shape.shape_type == 17:
            continue
        # Skip copying picture shapes like the Graphic 24 warning triangle!
        if shape.shape_type == 13:
            continue
            
        el = shape.element
        new_el = copy.deepcopy(el)
        dest_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')


def clear_unused_placeholders(slide):
    """Remove empty template placeholders to prevent visual clutter and prompt texts."""
    for shape in list(slide.shapes):
        if shape.is_placeholder and shape.placeholder_format.type in (2, 7, 10, 14): # Body text placeholders
            slide.shapes.element.remove(shape.element)


def add_diagram_box(slide, text_lines, left, top, width, height, fill_color, text_colors=None):
    """Create a beautifully formatted diagram rounded rectangle with center-aligned metadata text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background() # Clear border lines for modern look
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    
    # Clean padding margins
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    
    text_colors = text_colors or [RGBColor(255, 255, 255), RGBColor(226, 232, 240)]
    
    for idx, line in enumerate(text_lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(2) if idx > 0 else Pt(0)
        
        r = p.add_run()
        r.text = line
        r.font.name = "Segoe UI"
        
        if idx == 0:
            r.font.bold = True
            r.font.size = Pt(11.5)
            r.font.color.rgb = text_colors[0]
        else:
            r.font.bold = False
            r.font.size = Pt(9.5)
            r.font.color.rgb = text_colors[1]
            
    return shape


def build_architecture_diagram(slide, left, top, max_width, max_height):
    """Draw a spectacular, symmetrical 3-Tier block architecture diagram on the right half of Slide 7."""
    col_width = Inches(4.80)
    box_height = Inches(0.95)
    
    # Symmetrical center alignment of the diagram column: (max_width - col_width)/2
    x_offset = left + (max_width - col_width) / 2
    
    # Box colors derived from WorkFusion Premium palette
    clr_presentation = RGBColor(15, 23, 42)      # `#0F172A` (Slate Dark)
    clr_business = RGBColor(16, 185, 129)       # `#10B981` (Primary Emerald)
    clr_database = RGBColor(51, 65, 85)         # `#334155` (Slate Neutral)
    clr_ai_engine = RGBColor(5, 150, 105)        # `#059669` (Success Emerald)
    clr_arrows = RGBColor(148, 163, 184)         # `#94A3B8` (Soft Blue-Grey)
    
    # Text metadata
    meta_presentation = ["Presentation Layer (Next.js)", "React 18 • TypeScript • TailwindCSS • Framer Motion"]
    meta_business = ["Business Logic Layer (Node.js API)", "Express.js Orchestrator • Security Middleware • JWT & bcrypt"]
    meta_database = ["Database (MongoDB Atlas)", "Cloud Document Store • Search Indexes"]
    meta_ai_engine = ["AI Recommendation (Python)", "FastAPI • scikit-learn • TF-IDF Cosine"]
    
    # ------------------ Draw Boxes ------------------
    # Tier 1 Box (Presentation Layer)
    y1 = top
    add_diagram_box(slide, meta_presentation, x_offset, y1, col_width, box_height, clr_presentation)
    
    # Tier 2 Box (Business Logic Layer)
    y2 = top + Inches(1.40)
    add_diagram_box(slide, meta_business, x_offset, y2, col_width, box_height, clr_business)
    
    # Tier 3 (Split Databases & Recommendation Service Symmetrically)
    y3 = top + Inches(2.80)
    box3_width = Inches(2.30)
    gap3_width = Inches(0.20)
    
    # Left database box
    add_diagram_box(slide, meta_database, x_offset, y3, box3_width, box_height, clr_database,
                     text_colors=[RGBColor(255, 255, 255), RGBColor(203, 213, 225)])
    
    # Right AI box
    add_diagram_box(slide, meta_ai_engine, x_offset + box3_width + gap3_width, y3, box3_width, box_height, clr_ai_engine,
                     text_colors=[RGBColor(255, 255, 255), RGBColor(167, 243, 208)])
    
    # ------------------ Draw Connecting Arrows ------------------
    arrow_w = Inches(0.30)
    arrow_h = Inches(0.30)
    
    # Arrow 1 (Presentation -> Business Logic)
    arrow1_left = x_offset + (col_width - arrow_w) / 2
    arrow1_top = y1 + box_height + Inches(0.07)
    ar1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, arrow1_left, arrow1_top, arrow_w, arrow_h)
    ar1.fill.solid()
    ar1.fill.fore_color.rgb = clr_arrows
    ar1.line.fill.background()
    
    # Arrow 2 (Business Logic -> MongoDB Database)
    arrow2_left = x_offset + (box3_width - arrow_w) / 2
    arrow2_top = y2 + box_height + Inches(0.07)
    ar2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, arrow2_left, arrow2_top, arrow_w, arrow_h)
    ar2.fill.solid()
    ar2.fill.fore_color.rgb = clr_arrows
    ar2.line.fill.background()
    
    # Arrow 3 (Business Logic -> FastAPI AI Service)
    arrow3_left = x_offset + box3_width + gap3_width + (box3_width - arrow_w) / 2
    arrow3_top = y2 + box_height + Inches(0.07)
    ar3 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, arrow3_left, arrow3_top, arrow_w, arrow_h)
    ar3.fill.solid()
    ar3.fill.fore_color.rgb = clr_arrows
    ar3.line.fill.background()


def update_slide_numbers_and_footers(prs):
    """Programmatically sets sequential page numbers and clean footers on all slides (including added ones)."""
    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue  # Skip title slide
            
        has_num = False
        has_foot = False
        
        # Check and update existing placeholders or textboxes representing slide number/footer
        for shape in list(slide.shapes):
            name_l = shape.name.lower()
            text_val = shape.text_frame.text.strip() if shape.has_text_frame else ""
            
            # Detect if it represents a slide number (placeholder or custom text box)
            is_num_shape = (
                "slide number" in name_l 
                or "slide_number" in name_l 
                or (shape.is_placeholder and shape.placeholder_format.type == 13)
                or (shape.shape_type == 17 and text_val.isdigit())
            )
            # Detect if it represents a footer (placeholder or custom footer text box)
            is_foot_shape = (
                "footer" in name_l 
                or (shape.is_placeholder and shape.placeholder_format.type == 15)
                or (shape.shape_type == 17 and "BSc." in text_val)
            )
            
            if is_num_shape:
                if shape.has_text_frame:
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.RIGHT
                    r = p.add_run()
                    r.text = str(i + 1)
                    r.font.name = "Segoe UI"
                    r.font.size = Pt(11)
                    r.font.color.rgb = RGBColor(128, 128, 128)
                    has_num = True
            elif is_foot_shape:
                if shape.has_text_frame:
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    r = p.add_run()
                    r.text = "BSc. Information Engg. Technology"
                    r.font.name = "Segoe UI"
                    r.font.size = Pt(9.5)
                    r.font.color.rgb = RGBColor(128, 128, 128)
                    has_foot = True
                    
        # Add custom beautifully styled footer and page numbers for slides that lack them (newly added slides)
        if not has_num:
            num_box = slide.shapes.add_textbox(Inches(10.83), Inches(7.07), Inches(1.44), Inches(0.4))
            tf_num = num_box.text_frame
            p_num = tf_num.paragraphs[0]
            p_num.alignment = PP_ALIGN.RIGHT
            r_num = p_num.add_run()
            r_num.text = str(i + 1)
            r_num.font.name = "Segoe UI"
            r_num.font.size = Pt(11)
            r_num.font.color.rgb = RGBColor(128, 128, 128)
            
        if not has_foot:
            foot_box = slide.shapes.add_textbox(Inches(4.03), Inches(7.07), Inches(5.27), Inches(0.4))
            tf_foot = foot_box.text_frame
            p_foot = tf_foot.paragraphs[0]
            p_foot.alignment = PP_ALIGN.CENTER
            r_foot = p_foot.add_run()
            r_foot.text = "BSc. Information Engg. Technology"
            r_foot.font.name = "Segoe UI"
            r_foot.font.size = Pt(9.5)
            r_foot.font.color.rgb = RGBColor(128, 128, 128)


def reorder_slides(path):
    """Place AI + Demo slides before Conclusion for a logical FYP flow in XML representation."""
    import zipfile
    import shutil

    # Reordered slide sequence:
    # 1. Title Slide (original 1)
    # 2. Problem Statement (original 2)
    # 3. Literature Review (original 3)
    # 4. Business Scope Section Cover (original 4)
    # 5. Business Scope & Platform Roles (new 12, added 1st)
    # 6. Methodology/Framework Section Cover (original 5)
    # 7. System Architecture & Methodology (new 13, added 2nd)
    # 8. Tools & Technology Stack (original 6)
    # 9. Constraints & Limitations (original 7)
    # 10. Work Breakdown & Timelines (original 8)
    # 11. AI Recommendation Engine (new 14, added 3rd)
    # 12. Employer Workspace (new 15, added 4th)
    # 13. Seeker Workspace (new 16, added 5th)
    # 14. Conclusion & Future Work (original 9)
    # 15. References (original 10)
    # 16. Any Questions? (original 11)
    order_files = [1, 2, 3, 4, 12, 5, 13, 6, 7, 8, 14, 15, 16, 9, 10, 11]
    tmp = path + ".tmp"

    with zipfile.ZipFile(path, "r") as zin:
        pres_xml = zin.read("ppt/presentation.xml").decode()
        rels_xml = zin.read("ppt/_rels/presentation.xml.rels").decode()

        rid_to_slide = {}
        for match in re.finditer(
            r'Id="(rId\d+)"[^>]+Target="slides/slide(\d+)\.xml"', rels_xml
        ):
            rid_to_slide[match.group(1)] = int(match.group(2))

        slide_to_rid = {value: key for key, value in rid_to_slide.items()}
        entries = re.findall(r'(<p:sldId[^>]+r:id="rId\d+"[^>]*/>)', pres_xml)

        new_entries = []
        for slide_num in order_files:
            rid = slide_to_rid[slide_num]
            for entry in entries:
                if f'r:id="{rid}"' in entry:
                    new_entries.append(entry)
                    break

        new_sldlst = "<p:sldIdLst>" + "".join(new_entries) + "</p:sldIdLst>"
        new_pres = re.sub(
            r"<p:sldIdLst>.*?</p:sldIdLst>",
            new_sldlst,
            pres_xml,
            flags=re.DOTALL,
        )

        with zipfile.ZipFile(tmp, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "ppt/presentation.xml":
                    data = new_pres.encode()
                zout.writestr(item, data)

    shutil.move(tmp, path)


def main():
    shutil.copy2(SRC, OUT)
    prs = Presentation(OUT)

    # =========================================================================
    # Slide 1 — Title (Standard layout, perfectly aligned)
    # =========================================================================
    slide = prs.slides[0]
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    if title_shape:
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "WorkFusion: One Platform For Every Skill"
        r.font.bold = True
        r.font.name = "Segoe UI"
        r.font.size = Pt(36)
        
    if subtitle_shape:
        tf = subtitle_shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.space_after = Pt(12)
        r = p.add_run()
        r.text = "AI-Powered Hybrid Employment Marketplace\n"
        r.font.bold = True
        r.font.name = "Segoe UI"
        r.font.size = Pt(18)
        
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = "Muhammad Sharjeel bin Riaz (F22-BS-CS-016)  |  Hamza Shahzad (F22-BS-CS-048)\n\n" \
                  "Live Demo: work-fusion2026-client-y2f9.vercel.app"
        r2.font.name = "Segoe UI"
        r2.font.size = Pt(14)

    # =========================================================================
    # Slide 2 — Problem Statement (Standard layout, perfectly aligned)
    # =========================================================================
    slide = prs.slides[1]
    populate_slide_title(slide, "Problem Statement")
    clear_unused_placeholders(slide)
    
    # Use exact coordinates of original body placeholder: left=1.20, top=2.02, width=7.06, height=4.40
    add_premium_bullets(
        slide,
        [
            "Market Fragmentation: Employers are forced to use separate portals (Upwork/Fiverr for online freelancing, Rozee.pk for corporate jobs, and informal word-of-mouth for physical services).",
            "Broken Recruitment Pipelines: Existing systems lack unified contract-based, structured hiring workflows spanning bidding, scheduling, and milestone releases.",
            "Opaque & Manual Matching: Standard platforms offer basic keyword filters rather than smart, explainable matching or automated skill gap feedback.",
            "Invisible Local Labor: Qualified physical service workers (electricians, plumbers) remain digitally isolated with no verified reputation index.",
            "Absence of Security Interlocks: Crucial communication gates, verified contractor reviews, and escrow milestone guards are completely missing."
        ],
        left=Inches(1.20), top=Inches(2.02), width=Inches(7.06), height=Inches(4.40),
        font_size_pt=14
    )
    # Clear overlay text in Rectangle 15
    for shape in slide.shapes:
        if shape.name == "Rectangle 15" and shape.has_text_frame:
            shape.text_frame.clear()

    # =========================================================================
    # Slide 3 — Literature Review (Style B: Sidebar split-screen layout)
    # =========================================================================
    slide = prs.slides[2]
    populate_slide_title(slide, "Literature Review")
    
    # Clean the table size and align it beautifully
    for shape in slide.shapes:
        if shape.has_table:
            shape.left = Inches(5.19)
            shape.top = Inches(1.20)
            shape.width = Inches(7.43)
            shape.height = Inches(3.20)
            format_table(shape.table, font_size_pt=10.5)

    # Remove any existing custom note textboxes to avoid overlaps
    for shape in list(slide.shapes):
        if shape.name == "WorkFusionNoteBox" or (shape.has_text_frame and "WorkFusion (Our Solution)" in shape.text_frame.text):
            slide.shapes.element.remove(shape.element)

    # Add a beautiful, styled WorkFusion callout textbox below the comparison table (aligned with right area: left=5.19, width=7.43)
    txbox = slide.shapes.add_textbox(Inches(5.19), Inches(4.80), Inches(7.43), Inches(1.50))
    txbox.name = "WorkFusionNoteBox"
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "WorkFusion (Our Unified Solution): "
    r1.font.bold = True
    r1.font.name = "Segoe UI"
    r1.font.size = Pt(12)
    
    r2 = p.add_run()
    r2.text = "Integrates digital freelancing and localized physical services under a strict, " \
              "interview-gated hiring pipeline with explainable 7-factor weighted AI recommendations " \
              "(achieving MAP = 1.00 / MRR = 1.00) and milestone feedback."
    r2.font.bold = False
    r2.font.name = "Segoe UI"
    r2.font.size = Pt(12)

    # =========================================================================
    # Slide 4 — Business Scope (Dark Section Cover Slide - Clean Cover Layout)
    # =========================================================================
    slide = prs.slides[3]
    populate_slide_title(slide, "Business Scope & Modules", font_size_pt=40)
    clear_unused_placeholders(slide)
    
    # Business Scope slide is a section cover (dark backdrop). Keep it perfectly minimal as designed.
    # No bullet points or screenshots on the cover slide; let it act as a beautiful visual divider!

    # =========================================================================
    # Slide 6 — Methodology/Framework (Dark Section Cover Slide - Clean Cover Layout)
    # =========================================================================
    slide = prs.slides[4]
    populate_slide_title(slide, "Methodology & Framework", font_size_pt=40)
    clear_unused_placeholders(slide)
    
    # Similarly, Methodology slide is a dark transition slide. Leave the body blank for elegant transition.

    # =========================================================================
    # Slide 8 — Tools & Technology Stack (Standard layout, perfectly aligned)
    # =========================================================================
    slide = prs.slides[5]
    populate_slide_title(slide, "Tools & Technology Stack")
    clear_unused_placeholders(slide)
    
    # Original placeholder: left=1.20in, top=2.29in, width=11.00in, height=4.14in
    add_premium_bullets(
        slide,
        [
            "Frontend Ecosystem: Next.js 14, React 18, TypeScript, TailwindCSS for responsive layouts, shadcn/ui components, and Framer Motion.",
            "Backend Orchestration: Node.js, Express.js REST APIs, Mongoose ODM, JSON Web Tokens (JWT), and bcrypt passwords.",
            "AI & Recommendation: Python 3.10, FastAPI, scikit-learn (for NLP, TF-IDF Vectorizer & Cosine Similarity), Pandas, and NumPy.",
            "Database & Hosting: Fully managed MongoDB Atlas, Vercel (frontend deployment), and Render (backend and FastAPI hosting).",
            "Development Tools: Git, npm, Playwright, and environments configured with secure dot-env variables."
        ],
        left=Inches(1.20), top=Inches(2.15), width=Inches(11.00), height=Inches(4.45),
        font_size_pt=14
    )

    # =========================================================================
    # Slide 9 — Constraints & Limitations (Style B: Sidebar layout)
    # =========================================================================
    slide = prs.slides[6]
    populate_slide_title(slide, "Constraints & Limitations")
    clear_unused_placeholders(slide)
    
    # Original Style B placeholder: left=5.19in, top=0.70in, width=7.43in, height=6.18in
    add_premium_bullets(
        slide,
        [
            "Simulated Escrow Gateway: Active financial payment gateways (like Stripe/EasyPaisa) are mocked; contract funding is currently simulated.",
            "Web-Only Platform: Built as a premium responsive web application; native iOS and Android packages are deferred to the future roadmap.",
            "Twin-Cities Calibration: Location proximity scoring and commute radiuses are primarily tailored for Islamabad and Rawalpindi.",
            "Communication Interlocks: Live text chat remains locked until the applicant reaches the 'Interview' stage to prevent recruiter spam.",
            "Stateless Recommendation Service: The Python FastAPI engine does not persist state, recalculating similarity scores dynamically on every request."
        ],
        left=Inches(5.19), top=Inches(0.70), width=Inches(7.43), height=Inches(6.18),
        font_size_pt=14,
        is_sidebar_slide=True
    )

    # =========================================================================
    # Slide 10 — Work Breakdown & Timelines (Standard layout, perfectly aligned)
    # =========================================================================
    slide = prs.slides[7]
    populate_slide_title(slide, "Work Breakdown & Timelines")
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            replacements = {
                "Backend Testing": ("Stateless AI Matching & NLP", "01-02-2026", "20-02-2026", "20", "Sharjeel"),
                "Frontend Testing": ("Admin Panels & Core Analytics", "15-02-2026", "28-02-2026", "14", "Hamza"),
                "Integration Testing": ("Production Cloud Deployment", "01-03-2026", "15-03-2026", "15", "Sharjeel & Hamza"),
                "Documentation Writing": ("Comprehensive Quality Assurance", "16-03-2026", "30-05-2026", "75", "Sharjeel & Hamza"),
                "Final Submission": ("Final Thesis & Presentation", "01-06-2026", "18-06-2026", "18", "Sharjeel & Hamza")
            }
            for row in table.rows:
                act = row.cells[1].text.strip()
                if act in replacements:
                    new_act, start, end, dur, resp = replacements[act]
                    row.cells[1].text = new_act
                    row.cells[2].text = start
                    row.cells[3].text = end
                    row.cells[4].text = dur
                    row.cells[5].text = resp
            format_table(table, font_size_pt=9.5)

    # =========================================================================
    # Slide 14 — Conclusion & Future Work (Style B: Sidebar layout)
    # =========================================================================
    slide = prs.slides[8]
    populate_slide_title(slide, "Conclusion & Future Work")
    clear_unused_placeholders(slide)
    
    # Original Style B placeholder: left=5.19in, top=0.70in, width=7.43in, height=6.18in
    add_premium_bullets(
        slide,
        [
            "Successful Unification: Developed Pakistan's premier hybrid marketplace combining online digital freelancing and local physical task outsourcing.",
            "Validated Recommendation Model: Created an explainable weighted AI recommendations microservice achieving high evaluation accuracy (MAP = 1.00, MRR = 1.00).",
            "Anti-Spam Gating: Enforced robust hiring workflow interlocks, ensuring chat remains closed until interviews and reviews gate completed projects.",
            "Startup-Grade Aesthetics: Crafted a highly responsive, modern interface inspired by Stripe, Vercel, and Apple guidelines.",
            "Future Horizon: Highly modular architecture prepared for real-time payment gateways, voice-interview rooms, and mobile Progressive Web Apps (PWA)."
        ],
        left=Inches(5.19), top=Inches(0.70), width=Inches(7.43), height=Inches(6.18),
        font_size_pt=14,
        is_sidebar_slide=True
    )

    # =========================================================================
    # Slide 15 — References (Style B: Sidebar layout)
    # =========================================================================
    slide = prs.slides[9]
    populate_slide_title(slide, "References")
    clear_unused_placeholders(slide)
    
    # Original Style B placeholder: left=5.19in, top=0.66in, width=7.01in, height=6.17in
    add_premium_bullets(
        slide,
        [
            "[1] Sommerville, Software Engineering, 10th ed., Pearson Education, 2016.",
            "[2] MongoDB Inc., MongoDB Manual. Available: https://www.mongodb.com/docs/ (Accessed Jan. 2026).",
            "[3] Node.js Foundation, Node.js Documentation. Available: https://nodejs.org/docs/ (Accessed Jan. 2026).",
            "[4] OpenJS Foundation, Express.js API Reference. Available: https://expressjs.com/ (Accessed Jan. 2026).",
            "[5] Meta Platforms Inc., React: A JavaScript Library. Available: https://react.dev/ (Accessed Jan. 2026).",
            "[6] Vercel Inc., Next.js Documentation. Available: https://nextjs.org/docs (Accessed Jan. 2026).",
            "[7] FastAPI, FastAPI Web Framework. Available: https://fastapi.tiangolo.com/ (Accessed Jan. 2026).",
            "[8] scikit-learn, Machine Learning Library. Available: https://scikit-learn.org/ (Accessed Jan. 2026)."
        ],
        left=Inches(5.19), top=Inches(0.66), width=Inches(7.01), height=Inches(6.17),
        font_size_pt=12,
        is_sidebar_slide=True
    )

    # =========================================================================
    # Slide 16 — Any Questions? (Standard layout, perfectly aligned cover)
    # =========================================================================
    slide = prs.slides[10]
    populate_slide_title(slide, "Any Questions?")
    clear_unused_placeholders(slide)
    
    # Center credentials beautifully below the centralized title in the dark theme
    add_premium_bullets(
        slide,
        [
            "Live Production Website: https://work-fusion2026-client-y2f9.vercel.app/",
            "Employer Credentials: employer1@workfusion.com  (Password: password123)",
            "Candidate Credentials: seeker1@workfusion.com  (Password: password123)",
            "Project Engineers: Muhammad Sharjeel bin Riaz & Hamza Shahzad",
            "Platform Tagline: One Platform For Every Skill"
        ],
        left=Inches(1.20), top=Inches(2.15), width=Inches(11.00), height=Inches(4.45),
        font_size_pt=16
    )

    # =========================================================================
    # NEW SLIDES — Detailed Light Contents (Style A layouts, cloned from Slide 2)
    # =========================================================================
    title_content = prs.slide_layouts[1]
    
    # ----------------- 5. Business Scope & Platform Roles -----------------
    scope_slide = prs.slides.add_slide(title_content)
    populate_slide_title(scope_slide, "Business Scope & Platform Roles")
    clear_unused_placeholders(scope_slide)
    
    clone_slide_layout_shapes(prs.slides[1], scope_slide)
    
    add_premium_bullets(
        scope_slide,
        [
            "Employer Workspace: Publish digital or local physical jobs, evaluate AI candidate rankings with gap analysis, schedule interviews, and rate completed contracts.",
            "Service Seeker Workspace: Build profiles, receive personalized AI recommendations, apply to local/remote gigs, and chat directly once interviewed.",
            "Administrative Hub: Manage users, moderate job categories, monitor platform analytics, and suspend/restore accounts safely.",
            "Twin-City Commute Engine: Matches localized physical service requests based on commute distance and availability.",
            "Verified Feedback Loops: Reviews and rating releases are strictly locked until a job is marked 'Completed' by the employer."
        ],
        left=Inches(1.20), top=Inches(2.15), width=Inches(5.20), height=Inches(4.45),
        font_size_pt=13
    )
    homepage_img = os.path.join(ASSETS, "homepage.png")
    add_screenshot_with_aspect_ratio(scope_slide, homepage_img, left=Inches(6.90), top=Inches(2.15), max_width=Inches(5.23), max_height=Inches(4.45))

    # ----------------- 7. System Architecture & Methodology -----------------
    arch_slide = prs.slides.add_slide(title_content)
    populate_slide_title(arch_slide, "System Architecture & Methodology")
    clear_unused_placeholders(arch_slide)
    
    clone_slide_layout_shapes(prs.slides[1], arch_slide)
    
    add_premium_bullets(
        arch_slide,
        [
            "Premium 3-Tier Stack: Next.js frontend, Node.js + Express API backend orchestrator, and stateless Python FastAPI AI matching microservice.",
            "Strict Architectural Contract: The frontend never communicates directly with MongoDB or the Python service; Node.js manages all transactions.",
            "Stateless Python Matching: FastAPI service receives clean text profiles via HTTP POST, computes similarities, and returns raw ranking vectors instantly.",
            "Structured Pipeline Workflow: Standard development lifecycle followed: Requirement Analysis -> Database Design -> API Controllers -> AI Models -> React Frontend -> Integration Testing.",
            "Production-Grade Security: Implements JWT token authorization with HTTP-only refresh cookies, bcrypt hashing, Helmet protection, and rate limiting."
        ],
        left=Inches(1.20), top=Inches(2.15), width=Inches(5.20), height=Inches(4.45),
        font_size_pt=13
    )
    # Programmatically build a gorgeous vector diagram on the right half of Slide 7!
    build_architecture_diagram(arch_slide, left=Inches(6.90), top=Inches(2.15), max_width=Inches(5.23), max_height=Inches(4.45))

    # ----------------- 11. AI Recommendation Engine -----------------
    ai_slide = prs.slides.add_slide(title_content)
    populate_slide_title(ai_slide, "AI Recommendation Engine")
    clear_unused_placeholders(ai_slide)
    
    clone_slide_layout_shapes(prs.slides[1], ai_slide)
    
    add_premium_bullets(
        ai_slide,
        [
            "Advanced NLP Cleaners: Normalizes profiles, portfolio summaries, and job descriptions with hyphen preservation.",
            "TF-IDF Vectorization: Text properties are vectorized into numerical features using scikit-learn's stateless analyzer.",
            "7-Factor Weighted Scoring: Scores profiles on skill overlap (cosine similarity), experience, portfolio fit, rating, category preference, availability, and twin-city proximity.",
            "Skill Gap Analysis: Highlights missing skills (e.g. React, Docker) dynamically with recommendations.",
            "Explainable AI: Matches include score explanations (e.g., 'React matched, Portfolio matched, Commute inside 5km').",
            "Evaluation Metrics: Achieved a Mean Average Precision (MAP) of 1.00 and Mean Reciprocal Rank (MRR) of 1.00."
        ],
        left=Inches(1.20), top=Inches(2.15), width=Inches(5.20), height=Inches(4.45),
        font_size_pt=13
    )
    rec_img = os.path.join(ASSETS, "seeker-dashboard.png")
    add_screenshot_with_aspect_ratio(ai_slide, rec_img, left=Inches(6.90), top=Inches(2.15), max_width=Inches(5.23), max_height=Inches(4.45))

    # ----------------- 12. Employer Workspace -----------------
    emp_slide = prs.slides.add_slide(title_content)
    populate_slide_title(emp_slide, "Employer Workspace")
    clear_unused_placeholders(emp_slide)
    
    clone_slide_layout_shapes(prs.slides[1], emp_slide)
    
    add_premium_bullets(
        emp_slide,
        [
            "Clean Dashboard: Real-time overview of active jobs, applicant counts, scheduled interviews, and recent platform logs.",
            "Job Advertising: Create digital freelancing contracts or localized on-site physical service ads with specific skill and budget bounds.",
            "AI Candidate Ranking: Instantly ranks all job applicants as soon as the employer opens the job applicant details.",
            "Workflow Status Control: Move candidates from Pending to Reviewed, schedule Interviews, hire, or complete projects.",
            "Interview Chat Gating: The messaging module is dynamically unlocked ONLY when the candidate reaches the 'Interview' stage."
        ],
        left=Inches(1.20), top=Inches(2.15), width=Inches(5.20), height=Inches(4.45),
        font_size_pt=13
    )
    emp_img = os.path.join(ASSETS, "employer-dashboard.png")
    add_screenshot_with_aspect_ratio(emp_slide, emp_img, left=Inches(6.90), top=Inches(2.15), max_width=Inches(5.23), max_height=Inches(4.45))

    # ----------------- 13. Seeker Workspace -----------------
    seek_slide = prs.slides.add_slide(title_content)
    populate_slide_title(seek_slide, "Seeker Workspace")
    clear_unused_placeholders(seek_slide)
    
    clone_slide_layout_shapes(prs.slides[1], seek_slide)
    
    add_premium_bullets(
        seek_slide,
        [
            "Intelligent Dashboard: Instant list of top recommended matches complete with score explanations and gap details.",
            "Advanced Gigs Discovery: Browse active digital/physical listings with filters for category, budget, location, and skills.",
            "Interactive Timelines: Real-time visual timeline showing current application statuses and feedback progress.",
            "Direct Communication: Chat in real-time with employers once your application is moved to the 'Interview' stage.",
            "Profile Strengths Meter: Visual progress tracker displaying overall profile completion and identifying missing details."
        ],
        left=Inches(1.20), top=Inches(2.15), width=Inches(5.20), height=Inches(4.45),
        font_size_pt=13
    )
    seek_img = os.path.join(ASSETS, "seeker-explore.png")
    add_screenshot_with_aspect_ratio(seek_slide, seek_img, left=Inches(6.90), top=Inches(2.15), max_width=Inches(5.23), max_height=Inches(4.45))

    # =========================================================================
    # Final Presentation Adjustments
    # =========================================================================
    # Save the slides, reorder them to restore transition flows, reload, and apply numbers
    prs.save(OUT)
    reorder_slides(OUT)
    
    # Reload presentation to set sequential slide numbers correctly across the new order
    final_prs = Presentation(OUT)
    update_slide_numbers_and_footers(final_prs)
    final_prs.save(OUT)
    
    print(f"Updated premium presentation saved to {OUT} ({len(final_prs.slides)} slides)")


if __name__ == "__main__":
    main()
